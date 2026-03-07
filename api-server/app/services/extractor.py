from __future__ import annotations

import logging
import tempfile
from pathlib import Path

import fitz  # PyMuPDF — pure Python PDF renderer, no poppler needed
import pdfplumber
from PIL import Image

from app.config import settings
from app.models import ExtractionMethod, PageContent
from app.services.ocr import OCRService

logger = logging.getLogger(__name__)

# Optional: python-docx for Word documents
try:
    import docx as python_docx

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# Optional: python-pptx for PowerPoint
try:
    import pptx as python_pptx

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class DocumentExtractor:
    """Unified extractor for PDF, DOCX, PPTX, images, and plain text."""

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".txt", ".md", ".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"}

    def __init__(self, ocr_service: OCRService):
        self.ocr = ocr_service

    async def extract(self, file_path: Path, content_bytes: bytes | None = None) -> list[PageContent]:
        """Route to the correct extractor based on file extension."""
        ext = file_path.suffix.lower()

        if ext == ".pdf":
            return await self._extract_pdf(file_path)
        elif ext in (".docx", ".doc"):
            return self._extract_docx(file_path)
        elif ext == ".pptx":
            return self._extract_pptx(file_path)
        elif ext in (".txt", ".md"):
            return self._extract_text(file_path)
        elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"):
            return await self._extract_image(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    # ── PDF ───────────────────────────────────────────────────────────────────

    async def _extract_pdf(self, pdf_path: Path) -> list[PageContent]:
        pages: list[PageContent] = []
        native_texts: list[str] = []
        tables_per_page: list[str] = []

        with pdfplumber.open(str(pdf_path)) as pdf:
            total = len(pdf.pages)
            logger.info("PDF %s: %d pages", pdf_path.name, total)

            for page in pdf.pages:
                text = page.extract_text() or ""
                native_texts.append(text.strip())

                tables = page.extract_tables()
                tables_md = ""
                if tables:
                    for table in tables:
                        tables_md += _table_to_markdown(table) + "\n\n"
                tables_per_page.append(tables_md.strip())

        # Determine which pages need OCR
        needs_ocr = [
            i for i, text in enumerate(native_texts)
            if len(text) < settings.ocr_threshold or _is_garbled(text)
        ]

        ocr_texts: dict[int, str] = {}
        if needs_ocr:
            logger.info("OCR needed for pages: %s", [n + 1 for n in needs_ocr])
            pdf_doc = fitz.open(str(pdf_path))
            zoom = settings.dpi / 72
            matrix = fitz.Matrix(zoom, zoom)
            for idx in needs_ocr:
                pix = pdf_doc[idx].get_pixmap(matrix=matrix)
                img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                ocr_texts[idx] = await self.ocr.recognize_all(img)
            pdf_doc.close()

        for i in range(len(native_texts)):
            if i in ocr_texts:
                text = ocr_texts[i]
                method = ExtractionMethod.OCR
            else:
                text = native_texts[i]
                method = ExtractionMethod.NATIVE

            table_md = tables_per_page[i]
            has_tables = bool(table_md)
            if has_tables and method == ExtractionMethod.NATIVE:
                text = text + "\n\n" + table_md

            pages.append(PageContent(
                page_number=i + 1,
                text=text,
                extraction_method=method,
                has_tables=has_tables,
                tables_markdown=table_md,
            ))

        return pages

    # ── DOCX ──────────────────────────────────────────────────────────────────

    def _extract_docx(self, path: Path) -> list[PageContent]:
        if not HAS_DOCX:
            raise ImportError("python-docx is required for .docx files: pip install python-docx")

        doc = python_docx.Document(str(path))
        paragraphs: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                style = para.style.name if para.style else ""
                if "Heading" in style:
                    level = 1
                    for ch in style:
                        if ch.isdigit():
                            level = int(ch)
                            break
                    paragraphs.append(f"{'#' * level} {text}")
                else:
                    paragraphs.append(text)

        # Extract tables
        for table in doc.tables:
            paragraphs.append(_docx_table_to_markdown(table))

        full_text = "\n\n".join(paragraphs)
        return [PageContent(
            page_number=1,
            text=full_text,
            extraction_method=ExtractionMethod.NATIVE,
        )]

    # ── PPTX ──────────────────────────────────────────────────────────────────

    def _extract_pptx(self, path: Path) -> list[PageContent]:
        if not HAS_PPTX:
            raise ImportError("python-pptx is required for .pptx files: pip install python-pptx")

        prs = python_pptx.Presentation(str(path))
        pages: list[PageContent] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    if rows:
                        texts.append(_rows_to_markdown(rows))

            if texts:
                pages.append(PageContent(
                    page_number=slide_num,
                    text="\n\n".join(texts),
                    extraction_method=ExtractionMethod.NATIVE,
                ))

        return pages

    # ── Plain text ────────────────────────────────────────────────────────────

    def _extract_text(self, path: Path) -> list[PageContent]:
        text = path.read_text(encoding="utf-8", errors="replace")
        return [PageContent(
            page_number=1,
            text=text,
            extraction_method=ExtractionMethod.NATIVE,
        )]

    # ── Image ─────────────────────────────────────────────────────────────────

    async def _extract_image(self, path: Path) -> list[PageContent]:
        image = Image.open(str(path))
        text = await self.ocr.recognize_all(image)
        return [PageContent(
            page_number=1,
            text=text,
            extraction_method=ExtractionMethod.OCR,
        )]


# ── Helpers ──────────────────────────────────────────────────────────────────


def _table_to_markdown(table: list[list]) -> str:
    if not table or len(table) < 1:
        return ""
    cleaned = []
    for row in table:
        cleaned.append([
            str(cell).replace("\n", " ").strip() if cell else ""
            for cell in row
        ])
    return _rows_to_markdown(cleaned)


def _rows_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    header = rows[0]
    md = "| " + " | ".join(header) + " |\n"
    md += "| " + " | ".join(["---"] * len(header)) + " |\n"
    for row in rows[1:]:
        while len(row) < len(header):
            row.append("")
        md += "| " + " | ".join(row[: len(header)]) + " |\n"
    return md


def _docx_table_to_markdown(table) -> str:
    rows = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return _rows_to_markdown(rows)


def _is_garbled(text: str) -> bool:
    """Detect garbled text from PDF font encoding issues.

    Checks if the text has an abnormally low ratio of real dictionary words
    or common patterns. Garbled text like '6FDOLQJ)DFWRU' passes alphanumeric
    checks but fails linguistic coherence.
    """
    if len(text) < 50:
        return False

    words = text.split()
    if not words:
        return False

    # Check ratio of words that contain a vowel (real words almost always do)
    vowel_words = sum(
        1 for w in words
        if any(c in w.lower() for c in "aeiouy")
    )
    vowel_ratio = vowel_words / len(words)

    # Check ratio of words with excessive uppercase-to-length ratio
    # (garbled text often has patterns like "6FDOLQJ" — mostly uppercase with digits)
    garbled_words = sum(
        1 for w in words
        if len(w) > 4 and sum(1 for c in w if c.isupper()) / len(w) > 0.6
        and not w.isupper()  # all-caps acronyms are fine
    )
    garbled_ratio = garbled_words / len(words) if words else 0

    return vowel_ratio < 0.5 or garbled_ratio > 0.3
